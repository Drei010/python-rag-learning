from pathlib import Path
from typing import List, Tuple

from langchain_core.documents import Document

from core.config import settings
from services.file_service import get_supported_files


def get_powerpoint_files() -> List[Path]:
    return sorted(
        path
        for path in get_supported_files()
        if path.is_file()
        and path.suffix.lower() in settings.supported_powerpoint_extensions
    )


def iter_shape_text(shape) -> List[str]:
    texts = []

    if getattr(shape, "has_text_frame", False) and shape.text.strip():
        texts.append(shape.text.strip())

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    texts.append(text)

    return texts


def slide_to_document(file_path: Path, slide_number: int, slide) -> Document:
    texts = []

    for shape in slide.shapes:
        texts.extend(iter_shape_text(shape))

    page_content = "\n".join(
        [
            f"Source file: {file_path.name}",
            "File type: PowerPoint",
            f"Slide: {slide_number}",
            *texts,
        ]
    )

    return Document(
        page_content=page_content,
        metadata={
            "source": file_path.name,
            "file_type": "powerpoint",
            "slide": slide_number,
        },
        id=f"{file_path.name}:slide:{slide_number}",
    )


def load_powerpoint_documents() -> Tuple[List[Document], List[str]]:
    from pptx import Presentation

    loaded_documents = []
    loaded_ids = []

    for powerpoint_file in get_powerpoint_files():
        presentation = Presentation(str(powerpoint_file))

        for slide_index, slide in enumerate(presentation.slides, start=1):
            document = slide_to_document(powerpoint_file, slide_index, slide)
            if document.page_content.strip():
                loaded_documents.append(document)
                loaded_ids.append(document.id)

    return loaded_documents, loaded_ids
